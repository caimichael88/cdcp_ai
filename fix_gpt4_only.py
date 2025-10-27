#!/usr/bin/env python3
"""
Fix slides to say GPT-4 only (not Claude) for evaluation
"""

from pptx import Presentation

def fix_gpt4_references():
    """Update all slides to say GPT-4 only, not GPT-4/Claude"""

    prs = Presentation("CDCP_AI_Demo_Day_Complete.pptx")

    print("🔧 Fixing GPT-4 references (removing Claude)...\n")

    # Process all slides
    for i, slide in enumerate(prs.slides, 1):
        print(f"Checking Slide {i}...")
        for shape in slide.shapes:
            if hasattr(shape, "text_frame"):
                text = shape.text
                if "GPT-4/Claude" in text or "GPT-4 / Claude" in text:
                    # Replace with GPT-4 only
                    new_text = text.replace("GPT-4/Claude", "GPT-4")
                    new_text = new_text.replace("GPT-4 / Claude", "GPT-4")
                    shape.text_frame.clear()
                    p = shape.text_frame.paragraphs[0]
                    p.text = new_text
                    print(f"  ✓ Updated Slide {i}: GPT-4/Claude → GPT-4")

    # Save
    output_path = "CDCP_AI_Demo_Day_Complete.pptx"
    prs.save(output_path)
    print(f"\n✅ Fixed: All references now say 'GPT-4' only")

    return output_path

if __name__ == "__main__":
    try:
        output_file = fix_gpt4_references()
        print(f"\n🎉 Corrected!")
        print(f"   • Supervisor: GPT-4 (not Claude)")
        print(f"   • Fallback: GPT-4 (not Claude)")
        print(f"   • Code supports both, but default is GPT-4")
    except Exception as e:
        print(f"❌ Error: {e}")
        import traceback
        traceback.print_exc()
