#!/usr/bin/env python3
"""
Script to update the Demo Day presentation with CDCP AI project information
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import os

def update_presentation():
    """Update the Demo Day presentation with CDCP AI information"""

    # Load the existing presentation
    prs = Presentation("Demo Day Exaple Deck.pptx")

    print(f"Loaded presentation with {len(prs.slides)} slides")

    # Slide 1: Title slide (already has title and subtitle)
    # Already filled: CDCP AI, Voice Chatbox for Canada Dental Care Plan
    print("✓ Slide 1: Title slide already configured")

    # Slide 2: Problem
    slide2 = prs.slides[1]
    problem_text = """Canadians face significant barriers accessing dental care information:

• Complex eligibility requirements and coverage details
• Language barriers and literacy challenges
• Limited access to multilingual support resources
• Difficulty navigating government healthcare portals

Over 9 million Canadians lack dental coverage, and many eligible individuals don't understand CDCP benefits or how to access them."""

    update_slide_content(slide2, problem_text, "PROBLEM")
    print("✓ Slide 2: Problem updated")

    # Slide 3: Solution
    slide3 = prs.slides[2]
    solution_text = """CDCP AI is an intelligent voice chatbot that:

• Provides instant answers to CDCP questions via voice conversation
• Supports multilingual interactions with real-time speech-to-text
• Uses fine-tuned AI models trained on official CDCP documentation
• Delivers accurate, context-aware responses with RAG technology
• Converts responses back to natural speech for accessibility

Built with LangGraph agent orchestration, FastAPI backend, and React frontend."""

    update_slide_content(slide3, solution_text, "SOLUTION")
    print("✓ Slide 3: Solution updated")

    # Slide 4-8: Additional content
    # We'll add more details to other slides as needed

    # Save the updated presentation
    output_path = "CDCP_AI_Demo_Day_Presentation.pptx"
    prs.save(output_path)
    print(f"\n✅ Presentation saved as: {output_path}")

    return output_path

def update_slide_content(slide, content_text, title_check=None):
    """Update text content in a slide"""
    for shape in slide.shapes:
        if hasattr(shape, "text_frame"):
            text_frame = shape.text_frame
            # Check if this is a content area (not title, not number)
            if len(text_frame.paragraphs) > 2:  # Content areas typically have multiple paragraphs
                # Clear existing content but keep formatting
                for i, paragraph in enumerate(text_frame.paragraphs):
                    if i == 0:
                        paragraph.text = content_text
                    else:
                        paragraph.text = ""
                break

if __name__ == "__main__":
    try:
        output_file = update_presentation()
        print(f"\n🎉 Success! Your presentation is ready: {output_file}")
    except Exception as e:
        print(f"❌ Error: {e}")
        import traceback
        traceback.print_exc()
