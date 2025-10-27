#!/usr/bin/env python3
"""
Final comprehensive update for all slides with complete information
"""

from pptx import Presentation

def update_all_slides():
    """Update all slides with comprehensive, accurate information"""

    prs = Presentation("CDCP_AI_Demo_Day_Complete.pptx")

    print("📊 Updating all slides with complete information...\n")

    # SLIDE 3: Solution (High-level features)
    print("✍️ Slide 3: SOLUTION")
    slide3 = prs.slides[2]

    # Clear all text first
    for shape in slide3.shapes:
        if hasattr(shape, "text_frame"):
            shape.text_frame.clear()

    slide3_content = """SOLUTION

CDCP AI: Voice-enabled chatbot for instant, accurate CDCP answers

Key Features & Innovation:

• Multi-model AI architecture - Vector DB + Fine-tuned Llama-3.2-1B + GPT-4 evaluation and fallback

• Smart RAG pipeline - Semantic search retrieves context, trained Llama generates answers

• Dual-layer quality control - GPT-4 validates every answer; auto-fallback for out-of-scope or poor responses

• Voice-enabled & multilingual - Natural conversation interface, accessible to all Canadians

• Reduces dental clinic workload - Handles repetitive CDCP questions 24/7

Technology Stack:
Python FastAPI, LangChain, LangGraph, React/TypeScript, ChromaDB,
Llama-3.2-1B, GPT-4, Whisper ASR, Coqui TTS"""

    # Add to largest text box
    largest_shape = max(
        (s for s in slide3.shapes if hasattr(s, "text_frame")),
        key=lambda s: s.width * s.height,
        default=None
    )
    if largest_shape:
        p = largest_shape.text_frame.paragraphs[0]
        p.text = slide3_content

    # SLIDE 4: How It Works (Technical Architecture)
    print("✍️ Slide 4: HOW IT WORKS")
    slide4 = prs.slides[3]

    # Clear all text first
    for shape in slide4.shapes:
        if hasattr(shape, "text_frame"):
            shape.text_frame.clear()

    slide4_content = """HOW IT WORKS - Technical Architecture

End-to-End Voice Conversation Flow:
1. Voice Input → Whisper ASR → Text transcription
2. Query → ChromaDB Vector Search → Top 3 relevant CDCP documents
3. Context + Query → Fine-tuned Llama-3.2-1B-Instruct → Answer
4. Answer → Supervisor Evaluation (GPT-4/Claude)
5. If GOOD: Answer → Coqui TTS → Voice Output
   If BAD: GPT-4 Fallback → Coqui TTS → Voice Output

RAG Pipeline Components:

Vector Database (ChromaDB):
• Stores embeddings of official CDCP documentation (scraped from canada.ca)
• Sentence-Transformers model: all-MiniLM-L6-v2
• Semantic search retrieves top 3 most relevant document chunks
• Chunk size: 512 tokens with 50 token overlap

Fine-tuned Model (Llama-3.2-1B-Instruct):
• Base model: Meta's Llama-3.2-1B-Instruct
• Training: Multi-dialogue Q&A approach on CDCP dataset
• Generates context-aware, conversational answers
• Optimized for accessible, clear responses

Quality Control (Dual-Model):
• Supervisor LLM (GPT-4/Claude) evaluates every answer
• Checks: CDCP-related? Accurate? Helpful?
• Confidence scoring: HIGH/MEDIUM/LOW
• Intelligent fallback to GPT-4 if quality is poor or out-of-scope

Agent Orchestration (LangGraph):
• Multi-step workflow: ASR → RAG → Evaluation → TTS
• State management and tool routing
• Intelligent decision-making based on quality scores

Data Ingestion Process:
• Web scraping: Official canada.ca CDCP pages
• Document chunking: 512 tokens per chunk
• Embedding generation: Sentence-Transformers
• Vector storage: ChromaDB persistent database

Technology Stack:
• Agent: LangChain, LangGraph (orchestration & workflow)
• Backend: Python, FastAPI
• AI Models: Llama-3.2-1B (fine-tuned), GPT-4, Claude, Whisper
• Voice: OpenAI Whisper (ASR), Coqui TTS (speech synthesis)
• Database: ChromaDB vector store
• Frontend: React 19, TypeScript, Webpack, Nx monorepo"""

    largest_shape = max(
        (s for s in slide4.shapes if hasattr(s, "text_frame")),
        key=lambda s: s.width * s.height,
        default=None
    )
    if largest_shape:
        p = largest_shape.text_frame.paragraphs[0]
        p.text = slide4_content

    # SLIDE 5: Impact & Results
    print("✍️ Slide 5: IMPACT & RESULTS")
    slide5 = prs.slides[4]

    for shape in slide5.shapes:
        if hasattr(shape, "text_frame"):
            shape.text_frame.clear()

    slide5_content = """IMPACT & RESULTS

Accessibility Improvements:
• Voice-first interface removes literacy barriers
• Multilingual support for diverse Canadian population
• 24/7 availability - no wait times
• Reduces dental clinic staff workload on routine questions

Technical Achievements:
• Successfully scraped and ingested official CDCP documentation
• Built end-to-end RAG pipeline with semantic search
• Fine-tuned Llama-3.2-1B on CDCP Q&A dataset
• Implemented dual-model quality control system
• Integrated LangGraph for intelligent agent orchestration
• Complete voice conversation pipeline (ASR → AI → TTS)

System Performance:
• Semantic search retrieves relevant context accurately
• Quality control ensures reliable, on-topic responses
• Fallback mechanism handles out-of-scope questions gracefully
• Real-time voice processing with acceptable latency

Value Proposition:
• Dental clinics: Reduced time spent on repetitive CDCP inquiries
• Canadians: Easy access to CDCP information regardless of literacy/language
• Government: Scalable solution for public service information
• Future: Can be adapted for other government programs

Future Potential:
• Scale to provincial dental programs
• Mobile app deployment for wider accessibility
• Integration with official CDCP portal
• Multi-language expansion beyond English
• Voice authentication for personalized queries"""

    largest_shape = max(
        (s for s in slide5.shapes if hasattr(s, "text_frame")),
        key=lambda s: s.width * s.height,
        default=None
    )
    if largest_shape:
        p = largest_shape.text_frame.paragraphs[0]
        p.text = slide5_content

    # SLIDE 6: Challenges & Learnings
    print("✍️ Slide 6: CHALLENGES & LEARNINGS")
    slide6 = prs.slides[5]

    for shape in slide6.shapes:
        if hasattr(shape, "text_frame"):
            shape.text_frame.clear()

    slide6_content = """CHALLENGES & LEARNINGS

Technical Challenges:

Data Collection & Preparation:
• Web scraping official CDCP documentation with proper structure
• Chunking strategy to maintain context (512 tokens, 50 overlap)
• Creating high-quality Q&A dataset for fine-tuning

Model Fine-Tuning:
• Fine-tuning Llama-3.2-1B with limited compute resources
• Multi-dialogue training approach for conversational responses
• Balancing model size vs. accuracy vs. inference speed

RAG Pipeline:
• Selecting optimal embedding model (Sentence-Transformers)
• Determining right number of context documents (settled on top 3)
• Managing retrieval-generation balance

Quality Control:
• Implementing reliable evaluation metrics
• Supervisor LLM prompt engineering for consistent evaluation
• Handling edge cases and out-of-scope queries

Voice Processing:
• Managing audio processing latency
• Whisper ASR accuracy with different accents
• Coqui TTS voice quality and naturalness
• End-to-end voice pipeline integration

Key Learnings:

✓ RAG significantly improves answer accuracy over standalone fine-tuned model
✓ Dual-model architecture provides production-grade reliability
✓ LangGraph offers flexible, maintainable agent orchestration
✓ Quality control is essential for healthcare information systems
✓ Voice interfaces require careful UX design for accessibility
✓ Iterative testing and refinement crucial for performance
✓ Multi-model approach balances accuracy, cost, and latency

Solution Approaches:
• Comprehensive error logging for debugging
• Incremental development: RAG → Fine-tuning → Voice → Quality Control
• Modular architecture for easy testing and iteration
• Supervisor evaluation prevents hallucinations and off-topic responses"""

    largest_shape = max(
        (s for s in slide6.shapes if hasattr(s, "text_frame")),
        key=lambda s: s.width * s.height,
        default=None
    )
    if largest_shape:
        p = largest_shape.text_frame.paragraphs[0]
        p.text = slide6_content

    # SLIDE 7: Development Timeline
    print("✍️ Slide 7: DEVELOPMENT TIMELINE")
    slide7 = prs.slides[6]

    for shape in slide7.shapes:
        if hasattr(shape, "text_frame"):
            shape.text_frame.clear()

    slide7_content = """DEVELOPMENT TIMELINE & PROCESS

Phase 1: Research & Planning (Week 1-2)
• CDCP documentation analysis
• Technology stack selection (LangChain, LangGraph, FastAPI)
• RAG vs. Fine-tuning vs. Hybrid approach evaluation
• Architecture design (multi-model approach)

Phase 2: Data Collection & RAG Pipeline (Week 3-4)
• Web scraping official CDCP documentation
• Document chunking and embedding generation
• ChromaDB vector database setup
• RAG pipeline implementation and testing
• Semantic search optimization

Phase 3: Model Fine-Tuning (Week 4-5)
• Q&A dataset preparation from CDCP documentation
• Llama-3.2-1B fine-tuning with multi-dialogue approach
• Model evaluation and optimization
• Integration with RAG pipeline

Phase 4: Quality Control System (Week 5-6)
• Supervisor LLM implementation (GPT-4/Claude)
• Evaluation prompt engineering
• Dual-model architecture integration
• Fallback mechanism for poor/out-of-scope answers

Phase 5: Voice Integration (Week 6-7)
• OpenAI Whisper ASR integration
• Coqui TTS setup and configuration
• LangGraph agent orchestration
• End-to-end voice pipeline testing

Phase 6: Frontend & Integration (Week 7-8)
• React voice chatbot interface
• Real-time audio recording and playback
• API integration with Python backend
• User experience refinement

Phase 7: Testing & Refinement (Week 8-9)
• End-to-end system testing
• Performance optimization
• Edge case handling
• Demo preparation and documentation

Team: Michael Cai (Solo Project)
Tools: VS Code, GitHub, Nx monorepo, Python virtual environments"""

    largest_shape = max(
        (s for s in slide7.shapes if hasattr(s, "text_frame")),
        key=lambda s: s.width * s.height,
        default=None
    )
    if largest_shape:
        p = largest_shape.text_frame.paragraphs[0]
        p.text = slide7_content

    # SLIDE 8: Thank You / Q&A
    print("✍️ Slide 8: THANK YOU")
    slide8 = prs.slides[7]

    for shape in slide8.shapes:
        if hasattr(shape, "text_frame"):
            shape.text_frame.clear()

    slide8_content = """THANK YOU

Questions?

Project Summary:
CDCP AI - Voice-enabled chatbot with multi-model architecture
Using RAG + Fine-tuned Llama + GPT-4 quality control

Key Technologies:
• LangChain & LangGraph (agent orchestration)
• Llama-3.2-1B-Instruct (fine-tuned)
• ChromaDB (vector database)
• Whisper ASR + Coqui TTS (voice)
• FastAPI + React (full-stack)

Project Links:
• GitHub: [Your GitHub URL]
• Demo: [Live Demo URL if available]

Contact:
Michael Cai
Email: [Your Email]

Special Thanks:
• AI Class Instructors & TAs
• Canada.ca for CDCP documentation
• Open Source Community (LangChain, Hugging Face, Meta)
• Coqui TTS, OpenAI Whisper teams

Future Work:
• Production deployment
• Mobile app version
• Expand to other government services
• Multi-language support"""

    largest_shape = max(
        (s for s in slide8.shapes if hasattr(s, "text_frame")),
        key=lambda s: s.width * s.height,
        default=None
    )
    if largest_shape:
        p = largest_shape.text_frame.paragraphs[0]
        p.text = slide8_content

    # Save
    output_path = "CDCP_AI_Demo_Day_Complete.pptx"
    prs.save(output_path)
    print(f"\n✅ All slides updated: {output_path}")

    return output_path

if __name__ == "__main__":
    try:
        output_file = update_all_slides()
        print(f"\n🎉 Complete presentation ready!")
        print(f"\n📋 Slide Summary:")
        print(f"   Slide 1: Title - CDCP AI")
        print(f"   Slide 2: Problem - Dental clinic workload & info access")
        print(f"   Slide 3: Solution - Multi-model features & benefits")
        print(f"   Slide 4: How It Works - Complete RAG + fine-tuning architecture")
        print(f"   Slide 5: Impact - Achievements & value proposition")
        print(f"   Slide 6: Challenges - Technical hurdles & learnings")
        print(f"   Slide 7: Timeline - 9-week development process")
        print(f"   Slide 8: Thank You - Q&A & future work")
        print(f"\n📍 Next: Open in PowerPoint and add screenshots/images!")
    except Exception as e:
        print(f"❌ Error: {e}")
        import traceback
        traceback.print_exc()
