#!/usr/bin/env python3
"""
Reorganize Slide 3 (Solution - high level) and Slide 4 (How It Works - technical)
"""

from pptx import Presentation

def reorganize_slides():
    """Reorganize slides 3 and 4 for better presentation flow"""

    # Load the presentation
    prs = Presentation("CDCP_AI_Demo_Day_Complete.pptx")

    # Get Slide 3 (Solution - high level)
    slide3 = prs.slides[2]

    # Slide 3: Solution - High-level features and benefits
    slide3_content = """CDCP AI: An intelligent voice-enabled chatbot providing instant, accurate answers about the Canadian Dental Care Plan

What your solution does in one clear sentence:
Voice-to-voice AI assistant that helps Canadians and dental clinics get instant, accurate answers about CDCP eligibility, coverage, and benefits.

Key features:
• Voice-first interface - Talk naturally, get spoken responses
• Multilingual support for diverse Canadian population
• Instant answers to CDCP questions (eligibility, coverage, benefits)
• Accessible 24/7 - Reduces dental clinic workload
• Accurate responses powered by AI trained on official CDCP documentation

Why it matters:
• Reduces dental clinic staff time spent on repetitive CDCP questions
• Makes CDCP information accessible to all Canadians, regardless of literacy level
• Eliminates language barriers with multilingual voice support
• Provides reliable, verified answers to reduce confusion

Screenshot or visual of the product/interface:
[Voice chatbot interface with microphone, real-time conversation, and clear UI]

This is the WHAT and WHY - save technical details for next slide."""

    # Get Slide 4 (How It Works - technical architecture)
    slide4 = prs.slides[3]

    # Slide 4: How It Works - Technical deep-dive
    slide4_content = """HOW IT WORKS
Technical Architecture & RAG Pipeline

End-to-End Voice Conversation Flow:
1. Voice Input → ASR (OpenAI Whisper) → Text transcription
2. Text Query → Semantic Search (ChromaDB Vector Database)
3. Retrieved Context + Query → Fine-tuned Llama-3.2-1B-Instruct
4. Generated Answer → Quality Evaluation (GPT-4/Claude Supervisor)
5. If GOOD: Answer → TTS → Voice Output
   If BAD: Fallback to GPT-4 → TTS → Voice Output

RAG Pipeline Components:

Vector Database (ChromaDB):
• Stores embeddings of official CDCP documentation
• Semantic search retrieves top 3 most relevant documents
• Provides context for accurate answer generation

Fine-tuned Model (Llama-3.2-1B-Instruct):
• Trained with multi-dialogue approach on CDCP Q&A dataset
• Generates answers using query + retrieved context
• Optimized for conversational, accessible responses

Quality Control (Dual-Model Architecture):
• Supervisor LLM (GPT-4/Claude) evaluates every answer
• Checks: Is it CDCP-related? Accurate? Helpful?
• Confidence scoring: HIGH/MEDIUM/LOW
• Intelligent fallback to GPT-4 if quality is poor

Agent Orchestration (LangGraph):
• Manages multi-step workflow: ASR → RAG → Evaluation → TTS
• Intelligent routing based on quality scores
• Ensures reliable, production-ready responses

Technology Stack:
Backend: Python, FastAPI, LangChain, LangGraph
AI: Llama-3.2-1B (fine-tuned), GPT-4, Claude, Whisper
Database: ChromaDB (vector store)
Frontend: React, TypeScript, Axios"""

    # Update Slide 3
    print("Updating Slide 3 (Solution - High Level)...")
    updated_slide3 = False
    for shape in slide3.shapes:
        if hasattr(shape, "text_frame"):
            if "What yo" in shape.text or "solution does" in shape.text.lower() or "INCLUDE" in shape.text or "Key features" in shape.text or "CDCP AI" in shape.text:
                text_frame = shape.text_frame
                text_frame.clear()
                p = text_frame.paragraphs[0]
                p.text = slide3_content
                p.level = 0
                print("✓ Updated Slide 3: Solution (high-level benefits)")
                updated_slide3 = True
                break

    if not updated_slide3:
        print("⚠ Warning: Could not find Slide 3 content shape")

    # Update Slide 4
    print("\nUpdating Slide 4 (How It Works - Technical)...")
    updated_slide4 = False
    for shape in slide4.shapes:
        if hasattr(shape, "text_frame"):
            # Look for shapes with substantial text (likely the main content area)
            if len(shape.text) > 50:
                text_frame = shape.text_frame
                text_frame.clear()
                p = text_frame.paragraphs[0]
                p.text = slide4_content
                p.level = 0
                print("✓ Updated Slide 4: How It Works (technical architecture)")
                updated_slide4 = True
                break

    if not updated_slide4:
        print("⚠ Warning: Could not find Slide 4 content shape")

    # Save the updated presentation
    output_path = "CDCP_AI_Demo_Day_Complete.pptx"
    prs.save(output_path)
    print(f"\n✅ Presentation updated: {output_path}")

    return output_path

if __name__ == "__main__":
    try:
        output_file = reorganize_slides()
        print(f"\n🎉 Slides reorganized successfully!")
        print(f"\n📊 Slide Structure:")
        print(f"   Slide 3 (Solution): High-level features & benefits")
        print(f"      • What the chatbot does")
        print(f"      • Key user-facing features")
        print(f"      • Why it matters")
        print(f"\n   Slide 4 (How It Works): Technical architecture")
        print(f"      • RAG Pipeline flow diagram")
        print(f"      • Vector DB → Llama-3.2-1B → Supervisor")
        print(f"      • Dual-model quality control")
        print(f"      • Complete technology stack")
    except Exception as e:
        print(f"❌ Error: {e}")
        import traceback
        traceback.print_exc()
