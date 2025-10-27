#!/usr/bin/env python3
"""
Update Slide 2 (Problem) with corrected information
"""

from pptx import Presentation

def update_problem_slide():
    """Update the Problem slide with accurate information"""

    # Load the presentation
    prs = Presentation("CDCP_AI_Demo_Day_Complete.pptx")

    # Get Slide 2 (index 1)
    slide2 = prs.slides[1]

    # Updated problem content
    problem_content = """The problem or need you identified:
• Dental clinics are constantly receiving CDCP-related calls, adding significant extra workload
• Canadians have access to CDCP information but struggle to retrieve specific information they need
• Complex eligibility requirements and coverage details are difficult to navigate
• Language barriers make it harder for non-English speakers to understand their benefits

Who experiences this problem:
• Dental clinic staff overwhelmed with repetitive CDCP inquiries
• Canadian seniors, low-income families, and persons with disabilities
• Non-English speakers facing language barriers
• Anyone trying to understand their CDCP eligibility and coverage

Why it's important:
Dental clinics lose productive time answering routine questions, while eligible Canadians struggle to find clear answers about their coverage, leading to confusion, delayed care, and underutilization of available benefits."""

    # Find and update the text shape containing "INCLUDE"
    for shape in slide2.shapes:
        if hasattr(shape, "text_frame"):
            if "INCLUDE" in shape.text or "problem or need" in shape.text.lower():
                # Found the shape, update it
                text_frame = shape.text_frame
                text_frame.clear()
                p = text_frame.paragraphs[0]
                p.text = problem_content
                p.level = 0
                print("✓ Updated Problem slide content")
                break

    # Save the updated presentation
    output_path = "CDCP_AI_Demo_Day_Complete.pptx"
    prs.save(output_path)
    print(f"\n✅ Presentation updated: {output_path}")

    return output_path

if __name__ == "__main__":
    try:
        output_file = update_problem_slide()
        print(f"\n🎉 Slide 2 (Problem) has been updated with:")
        print(f"   • Dental clinic workload from CDCP calls")
        print(f"   • Information retrieval difficulties")
        print(f"   • Removed the '9+ million uninsured' statement")
    except Exception as e:
        print(f"❌ Error: {e}")
        import traceback
        traceback.print_exc()
