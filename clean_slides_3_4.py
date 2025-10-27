#!/usr/bin/env python3
"""
Clean slides 3 and 4 - remove ALL old text and add only new content
"""

from pptx import Presentation
from pptx.util import Pt

def clean_and_update_slides():
    """Clean all text from slides 3 and 4, then add fresh content"""

    # Load the presentation
    prs = Presentation("CDCP_AI_Demo_Day_Complete.pptx")

    # Get slides
    slide3 = prs.slides[2]  # Solution
    slide4 = prs.slides[3]  # How It Works

    print("🧹 Cleaning Slide 3 (Solution)...")
    # Clear ALL text from Slide 3
    for shape in slide3.shapes:
        if hasattr(shape, "text_frame"):
            text_frame = shape.text_frame
            text_frame.clear()

    print("🧹 Cleaning Slide 4 (How It Works)...")
    # Clear ALL text from Slide 4
    for shape in slide4.shapes:
        if hasattr(shape, "text_frame"):
            text_frame = shape.text_frame
            text_frame.clear()

    print("\n✍️ Adding new content to Slide 3...")
    # Add fresh content to Slide 3 - Solution (High-level)
    slide3_content = """SOLUTION

CDCP AI: Voice-enabled chatbot for instant, accurate CDCP answers

Key Features:
• Voice-first interface - talk naturally, get spoken responses
• Multilingual support for all Canadians
• 24/7 availability - reduces dental clinic workload
• AI-powered accurate responses from official CDCP documentation

Why It Matters:
• Reduces dental clinic staff time on repetitive questions
• Makes CDCP accessible regardless of literacy level
• Eliminates language barriers
• Provides reliable, verified information"""

    # Find a suitable text box to add content (usually the largest one)
    largest_shape = None
    largest_size = 0
    for shape in slide3.shapes:
        if hasattr(shape, "text_frame"):
            size = shape.width * shape.height
            if size > largest_size:
                largest_size = size
                largest_shape = shape

    if largest_shape:
        text_frame = largest_shape.text_frame
        p = text_frame.paragraphs[0]
        p.text = slide3_content
        print("✓ Slide 3 content added")
    else:
        print("⚠ Warning: Could not find text box in Slide 3")

    print("\n✍️ Adding new content to Slide 4...")
    # Add fresh content to Slide 4 - How It Works (Technical)
    slide4_content = """HOW IT WORKS
Technical Architecture

End-to-End Flow:
1. Voice Input → Whisper ASR → Text transcription
2. Query → ChromaDB Vector Search → Top 3 relevant documents
3. Context + Query → Fine-tuned Llama-3.2-1B-Instruct
4. Answer → Supervisor Evaluation (GPT-4/Claude)
5. If GOOD: TTS → Voice Output
   If BAD: GPT-4 Fallback → TTS → Voice Output

RAG Pipeline:
• Vector DB (ChromaDB): Semantic search of CDCP documentation
• Fine-tuned Llama-3.2-1B: Multi-dialogue training, context-aware generation
• Quality Control: Dual-model architecture with supervisor validation
• LangGraph: Intelligent agent orchestration

Technology Stack:
• Backend: Python, FastAPI, LangChain, LangGraph
• AI Models: Llama-3.2-1B (fine-tuned), GPT-4, Claude, Whisper
• Database: ChromaDB vector store
• Frontend: React, TypeScript"""

    # Find a suitable text box for Slide 4
    largest_shape = None
    largest_size = 0
    for shape in slide4.shapes:
        if hasattr(shape, "text_frame"):
            size = shape.width * shape.height
            if size > largest_size:
                largest_size = size
                largest_shape = shape

    if largest_shape:
        text_frame = largest_shape.text_frame
        p = text_frame.paragraphs[0]
        p.text = slide4_content
        print("✓ Slide 4 content added")
    else:
        print("⚠ Warning: Could not find text box in Slide 4")

    # Save
    output_path = "CDCP_AI_Demo_Day_Complete.pptx"
    prs.save(output_path)
    print(f"\n✅ Presentation cleaned and updated: {output_path}")

    return output_path

if __name__ == "__main__":
    try:
        output_file = clean_and_update_slides()
        print(f"\n🎉 Slides 3 & 4 cleaned successfully!")
        print(f"   • All old text removed")
        print(f"   • Fresh content added")
        print(f"   • No overlapping text")
    except Exception as e:
        print(f"❌ Error: {e}")
        import traceback
        traceback.print_exc()
