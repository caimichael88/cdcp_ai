#!/usr/bin/env python3
"""
Update Slide 3 (Solution) to include Vector DB, Llama model, and RAG pipeline details
"""

from pptx import Presentation

def update_solution_slide():
    """Update the Solution slide with complete technical architecture"""

    # Load the presentation
    prs = Presentation("CDCP_AI_Demo_Day_Complete.pptx")

    # Get Slide 3 (index 2)
    slide3 = prs.slides[2]

    # Updated solution content with Vector DB and Llama model details
    solution_content = """CDCP AI: An intelligent voice-enabled chatbot providing instant, accurate answers about the Canadian Dental Care Plan

What your solution does in one clear sentence:
Voice-to-voice AI assistant that uses RAG pipeline with fine-tuned Llama model and quality control to answer CDCP questions accurately and accessibly.

Key features and innovations:

RAG Pipeline Architecture:
• Vector Database (ChromaDB) for semantic search of CDCP documentation
• Fine-tuned Llama-3.2-1B-Instruct model trained with multi-dialogue approach
• Retrieval-Augmented Generation: Query → Vector search → Context + Fine-tuned model → Answer

Voice & Accessibility:
• Automatic Speech Recognition (ASR) using OpenAI Whisper
• Text-to-Speech (TTS) for natural voice responses
• Multilingual support for diverse Canadian population

Quality Control & Intelligence:
• Dual-model architecture: Fine-tuned Llama generates, Supervisor LLM validates
• GPT-4/Claude evaluates every answer for accuracy and relevance
• Intelligent fallback: If answer quality is poor, uses GPT-4 as backup
• LangGraph agent orchestration for smart conversation flow

Screenshot or visual of the product/interface:
[Voice chatbot interface with microphone, real-time transcription, and conversation history]

Technology Stack: Python FastAPI, React TypeScript, ChromaDB, Llama-3.2-1B, Whisper ASR"""

    # Find and update the text shape containing solution content
    updated = False
    for shape in slide3.shapes:
        if hasattr(shape, "text_frame"):
            if "What yo" in shape.text or "solution does" in shape.text.lower() or "INCLUDE" in shape.text or "Key features" in shape.text:
                # Found the shape, update it
                text_frame = shape.text_frame
                text_frame.clear()
                p = text_frame.paragraphs[0]
                p.text = solution_content
                p.level = 0
                print("✓ Updated Solution slide with Vector DB and Llama model details")
                updated = True
                break

    if not updated:
        print("⚠ Warning: Could not find solution content shape")

    # Save the updated presentation
    output_path = "CDCP_AI_Demo_Day_Complete.pptx"
    prs.save(output_path)
    print(f"\n✅ Presentation updated: {output_path}")

    return output_path

if __name__ == "__main__":
    try:
        output_file = update_solution_slide()
        print(f"\n🎉 Slide 3 (Solution) updated with:")
        print(f"   ✓ Vector Database (ChromaDB) for semantic search")
        print(f"   ✓ Fine-tuned Llama-3.2-1B-Instruct model")
        print(f"   ✓ Multi-dialogue training approach")
        print(f"   ✓ Complete RAG pipeline architecture")
        print(f"   ✓ Dual-model quality control system")
        print(f"   ✓ Intelligent fallback mechanism")
    except Exception as e:
        print(f"❌ Error: {e}")
        import traceback
        traceback.print_exc()
